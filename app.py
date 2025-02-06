import streamlit as st
import logging
from bs4 import BeautifulSoup
from docx import Document
from docx.shared import RGBColor
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls
from latex2mathml.converter import convert
from pygments import highlight
from pygments.lexers import get_lexer_by_name
from pygments.formatters import HtmlFormatter
from fpdf import FPDF
from pptx import Presentation
import io

# Setup logging
logging.basicConfig(level=logging.INFO)

def convert_latex_to_omml(latex):
    try:
        mathml = convert(latex)
        logging.info(f"Converted LaTeX to MathML: {mathml[:100]}")  # Log first 100 characters of MathML
        return f'<m:oMath xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">{mathml}</m:oMath>'
    except Exception as e:
        logging.error(f"Error converting LaTeX: {e}")
        return f"[Math: {latex}]"

def html_to_docx(html_content):
    if not html_content.strip():
        logging.warning("Empty HTML content received for DOCX conversion.")
        return None
    
    logging.info(f"HTML Content Received: {html_content[:500]}")  # Log first 500 characters of HTML
    
    soup = BeautifulSoup(html_content, "html.parser")
    logging.info(f"Parsed HTML content: {soup.prettify()[:1000]}")  # Log parsed HTML to debug
    
    doc = Document()

    # Look for specific tags and add content to doc
    for element in soup.find_all(True):
        logging.info(f"Processing Element: {element.name}")  # Log each element being processed
        
        if element.name == "p":
            doc.add_paragraph(element.get_text())
        elif element.name in ["h1", "h2", "h3"]:
            doc.add_paragraph(element.get_text(), style=element.name.capitalize())
        elif element.name == "code":
            p = doc.add_paragraph()
            run = p.add_run(element.get_text())
            run.font.name = "Courier New"
        elif element.name == "pre":
            lexer = get_lexer_by_name("python", stripall=True)
            formatter = HtmlFormatter(style="colorful")
            highlighted_code = highlight(element.get_text(), lexer, formatter)
            doc.add_paragraph(highlighted_code)
        elif element.name == "table":
            rows = element.find_all("tr")
            if rows:
                cols = len(rows[0].find_all(["td", "th"]))
                table = doc.add_table(rows=0, cols=cols)
                for row in rows:
                    cells = row.find_all(["td", "th"])
                    row_cells = table.add_row().cells
                    for i, cell in enumerate(cells):
                        row_cells[i].text = cell.get_text()
        elif element.name == "span" and "math" in element.get("class", []):
            latex_code = element.get_text()
            math_omml = convert_latex_to_omml(latex_code)
            p = doc.add_paragraph()
            p._element.append(parse_xml(math_omml))
    
    # After the loop, log the length of the document content
    logging.info(f"Document created with {len(doc.paragraphs)} paragraphs")
    
    return doc

def html_to_pdf(html_content):
    if not html_content.strip():
        logging.warning("Empty HTML content received for PDF conversion.")
        return None
    
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, html_content)
    return pdf

def html_to_ppt(html_content):
    if not html_content.strip():
        logging.warning("Empty HTML content received for PPT conversion.")
        return None
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide.shapes.add_textbox(50, 50, 600, 400)
    tf = textbox.text_frame
    tf.text = html_content
    return prs

# Streamlit UI
st.title("ChatGPT Export Tool")
st.write("Convert ChatGPT outputs into DOCX, PDF, or PowerPoint.")

uploaded_file = st.file_uploader("Upload HTML File", type=["html"], accept_multiple_files=False)
input_text = st.text_area("Paste HTML Content")
export_format = st.radio("Select Export Format", ["DOCX", "PDF", "PowerPoint"])

if st.button("Convert & Download"):
    if uploaded_file:
        html_content = uploaded_file.read().decode("utf-8")
        logging.info(f"File Uploaded. HTML Content Length: {len(html_content)}")
        st.write(html_content)  # Debugging output
    elif input_text.strip():
        html_content = input_text.strip()
        logging.info(f"HTML Content Provided. Length: {len(html_content)}")
        st.write(html_content)  # Debugging output
    else:
        st.error("Please provide input!")
        st.stop()
    
    if export_format == "DOCX":
        doc = html_to_docx(html_content)
        if doc:
            doc_path = "output.docx"
            doc.save(doc_path)
            logging.info(f"DOCX file saved at {doc_path}")
            with open(doc_path, "rb") as f:
                st.download_button("Download DOCX", f, file_name="output.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        else:
            st.error("Failed to generate DOCX file.")
    elif export_format == "PDF":
        pdf = html_to_pdf(html_content)
        if pdf:
            pdf_output = io.BytesIO()
            pdf.output(pdf_output, dest='S')
            pdf_output.seek(0)
            st.download_button("Download PDF", pdf_output, file_name="output.pdf", mime="application/pdf")
        else:
            st.error("Failed to generate PDF file.")
    elif export_format == "PowerPoint":
        ppt = html_to_ppt(html_content)
        if ppt:
            ppt_path = "output.pptx"
            ppt.save(ppt_path)
            with open(ppt_path, "rb") as f:
                st.download_button("Download PowerPoint", f, file_name="output.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        else:
            st.error("Failed to generate PowerPoint file.")
